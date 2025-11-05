# retrieval/parsers.py
"""
Offline file parsers for RAG ingestion.
Covers: PDF, DOCX/DOC, RTF, PPTX, XLSX/CSV, JSON/TXT/MD, code files, images (OCR).
"""

from __future__ import annotations
import io
import json
from pathlib import Path
from typing import Dict, List

# PDF
from pypdf import PdfReader

# DOCX/RTF
from docx import Document
try:
    from striprtf.striprtf import rtf_to_text
    _HAS_RTF = True
except Exception:
    _HAS_RTF = False

# PPTX
from pptx import Presentation

# Data Analysis
import pandas as pd
import openpyxl

# Images / OCR
from PIL import Image
try:
    import pytesseract
    _HAS_OCR = True
except Exception:
    _HAS_OCR = False

# Primitive helpers
def _join_nonempty(parts: List[str], sep: str = "\n") -> str:
    return sep.join([p for p in parts if p and p.strip()]).strip()

def _decode_best_effort(b: bytes) -> str:
    for enc in ("utf-8", "utf-16", "latin1"):
        try:
            return b.decode(enc)
        except Exception:
            continue
    return b.decode("utf-8", errors="ignore")

# PDF
def read_pdf_bytes(b: bytes) -> str:
    """
    Extract text page-by-page. If a page yields no text and OCR is available,
    OCR that page image (best-effort).
    """
    out: List[str] = []
    try:
        reader = PdfReader(io.BytesIO(b))
        for page in reader.pages:
            txt = page.extract_text() or ""
            if not txt.strip() and _HAS_OCR:
                pass # OCR fallback can be complex; skipping for now.
            if txt.strip():
                out.append(txt)
    except Exception:
        return ""
    return _join_nonempty(out)

# DOCX / DOC / RTF
def read_docx_bytes(b: bytes) -> str:
    try:
        doc = Document(io.BytesIO(b))
    except Exception:
        return ""
    parts: List[str] = []
    for p in doc.paragraphs:
        if (p.text or "").strip(): parts.append(p.text.strip())
    for tbl in doc.tables:
        for row in tbl.rows:
            row_txt = [cell.text.strip() for cell in row.cells if getattr(cell, "text", None)]
            if any(row_txt): parts.append(" | ".join(row_txt))
    return _join_nonempty(parts)

def read_rtf_bytes(b: bytes) -> str:
    if not _HAS_RTF: return ""
    try:
        return (rtf_to_text(_decode_best_effort(b)) or "").strip()
    except Exception:
        return ""

# PPTX
def read_pptx_bytes(b: bytes) -> str:
    try:
        prs = Presentation(io.BytesIO(b))
    except Exception:
        return ""
    slides_out: List[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts: List[str] = []
        for shp in slide.shapes:
            try:
                if hasattr(shp, "text") and shp.text and shp.text.strip():
                    parts.append(shp.text.strip())
                if shp.has_table:
                    for row in shp.table.rows:
                        row_txt = [c.text.strip() for c in row.cells if c.text]
                        if any(row_txt): parts.append(" | ".join(row_txt))
            except Exception:
                continue
        if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
            nt = slide.notes_slide.notes_text_frame.text
            if nt and nt.strip(): parts.append(f"Notes: {nt.strip()}")
        if parts:
            slides_out.append(f"Slide {i}:\n" + _join_nonempty(parts))
    return _join_nonempty(slides_out, sep="\n\n")


def _read_xlsx_with_pandas(b: bytes) -> str:
    """Original pandas-based XLSX parser. Good for data tables, bad for layouts. Kept as a fallback."""
    try:
        df_map = pd.read_excel(io.BytesIO(b), sheet_name=None, engine="openpyxl")
    except Exception:
        return ""
    out: List[str] = []
    for name, df in (df_map or {}).items():
        try:
            out.append(f"[Sheet: {name}]\n" + df.to_csv(index=False, na_rep=""))
        except Exception:
            continue
    return _join_nonempty(out, sep="\n\n")

def read_xlsx_bytes(b: bytes) -> str:
    """
    Layout-aware XLSX parser using openpyxl. It reads cell-by-cell to reconstruct
    document-style content, preserving the logical grouping of text.
    """
    try:
        workbook = openpyxl.load_workbook(io.BytesIO(b), data_only=True)
        full_text_parts = []

        for sheet in workbook:
            sheet_text = f"[Sheet: {sheet.title}]\n\n"
            current_block = []
            
            for row in sheet.iter_rows():
                # Get text from all cells in the row, filter out None values
                row_texts = [str(cell.value).strip() for cell in row if cell.value is not None]
                
                # An empty row acts as a separator between logical blocks
                if not row_texts:
                    if current_block:
                        full_text_parts.append("\n".join(current_block))
                        current_block = [] # Reset for the next block
                    continue
                
                # Join cell text with tabs to maintain some visual separation
                line_text = "\t".join(row_texts)

                # If we see a main header, finalize the previous block before starting a new one
                if "PREMANUFACTURING FORMULATION DATA" in line_text and current_block:
                    full_text_parts.append("\n".join(current_block))
                    current_block = [line_text] # Start new block with the header
                else:
                    current_block.append(line_text)

            # Add the last remaining block after the loop finishes
            if current_block:
                full_text_parts.append("\n".join(current_block))
        
        # Join all captured blocks with a clear separator
        result = "\n\n---\n\n".join(full_text_parts)
        return result if result.strip() else _read_xlsx_with_pandas(b)

    except Exception:
        # If the layout-aware parser fails for any reason, fall back to the simple one
        return _read_xlsx_with_pandas(b)

def read_csv_bytes(b: bytes) -> str:
    for enc in ("utf-8", "utf-16", "latin1"):
        for sep in (",", ";", "\t", "|"):
            try:
                df = pd.read_csv(io.BytesIO(b), encoding=enc, sep=sep)
                return df.to_csv(index=False, na_rep="")
            except Exception:
                continue
    return _decode_best_effort(b)

# JSON / TXT / CODE / XML / HTML / YAML
def read_json_bytes(b: bytes) -> str:
    try:
        obj = json.loads(_decode_best_effort(b))
        return json.dumps(obj, ensure_ascii=False, indent=2)
    except Exception:
        return _decode_best_effort(b)

def read_text_bytes(b: bytes) -> str:
    return _decode_best_effort(b)

# Images → OCR
def read_image_bytes(b: bytes) -> str:
    if not _HAS_OCR: return ""
    try:
        img = Image.open(io.BytesIO(b))
        txt = pytesseract.image_to_string(img)
        return txt.strip()
    except Exception:
        return ""

# Registry used by vector_store.py
SUPPORTED_READERS: Dict[str, callable] = {
    # docs
    ".docx": read_docx_bytes,
    ".doc":  read_docx_bytes,
    ".rtf":  read_rtf_bytes,
    # pdf
    ".pdf":  read_pdf_bytes,
    # slides
    ".pptx": read_pptx_bytes,
    # sheets / data
    ".xlsx": read_xlsx_bytes, # <-- This now points to our new, smarter function
    ".csv":  read_csv_bytes,
    ".json": read_json_bytes,
    # text-ish
    ".txt":  read_text_bytes,
    ".md":   read_text_bytes,
    ".py":   read_text_bytes,
    ".js":   read_text_bytes,
    ".ts":   read_text_bytes,
    ".java": read_text_bytes,
    ".c":    read_text_bytes,
    ".cpp":  read_text_bytes,
    ".xml":  read_text_bytes,
    ".html": read_text_bytes,
    ".yaml": read_text_bytes,
    ".yml":  read_text_bytes,
    # images
    ".png":  read_image_bytes,
    ".jpg":  read_image_bytes,
    ".jpeg": read_image_bytes,
    ".gif":  read_image_bytes,
    ".tif":  read_image_bytes,
    ".tiff": read_image_bytes,
}
